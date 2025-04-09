from report.models import (
    ONAFormAnswered,
    FormSubsectionAnswered,
    FormSectionAnswered,
    QuestionAnswer,
    ONAFormDistribution
)
from collections import Counter
import matplotlib.pyplot as plt
from datetime import datetime
import seaborn as sns
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from django.core.mail import EmailMessage
from reportlab.lib.pagesizes import letter
from data_management.models import Evaluator
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from django.db.models import QuerySet
from reportlab.platypus import Image
from reportlab.platypus import Table, TableStyle, PageBreak

import os
import matplotlib.dates as mdates
from django.utils import timezone
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont




class GraphsGenerator:
    def plot_bar_plot(self, data, title):
        # Assuming `data` is a dictionary that returns a distribution by percentage
        data = self.awnser_distribution_by_percentage(data)

        # Define the colors and their corresponding categories
        color_map = {
            "Não conforme": "#E41A1C",  # Red
            "Parcial conforme": "#377EB8",  # Blue
            "Conforme": "#4DAF4A",  # Green
            "Supera": "#FF7F00"  # Orange
        }

        # Convert data to a DataFrame
        df = pd.DataFrame(list(data.items()), columns=["Status", "Percentage"])
        df["Status"] = df["Status"].str.capitalize()

        # Sort the DataFrame if needed (optional)
        df = df.set_index('Status').reindex(["Não conforme", "Parcial conforme", "Conforme", "Supera"]).reset_index()

        # Set up the plot
        plt.figure(figsize=(10, 6))

        # Use a list comprehension to assign colors based on the Status
        colors = [color_map[status] for status in df["Status"]]

        # Create the bar plot
        ax = sns.barplot(x="Status", y="Percentage", data=df, palette=colors, edgecolor="black")

        # Customize the plot
        plt.ylabel("Conformidade (%)")
        plt.xlabel("Respostas", labelpad=15)
        plt.title(title)

        plt.grid(True, axis="y", linestyle="--", alpha=0.7)

        # Set y-axis range
        plt.ylim(0, 100)

        # Annotate bars with percentages
        for p in ax.patches:
            height = p.get_height()
            if height == 0:
                continue
            ax.annotate(
                f"{p.get_height():.1f}%",
                (p.get_x() + p.get_width() / 2.0, p.get_height()),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )

        # Create custom legend
        handles = []
        for status, color in color_map.items():
            handles.append(
                plt.Line2D(
                    [0],
                    [0],
                    marker="o",
                    color="w",
                    markerfacecolor=color,
                    markersize=10,
                    label=status,
                )
            )

        # Display the legend with the correct order
        plt.legend(handles=handles, title="Legenda")

        # Adjust layout
        plt.tight_layout()

        # Save the figure to a buffer (for embedding in PowerPoint or other purposes)
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)
        return buf

    def plot_grouped_bar(self, data, title):
        # Assuming `data` is already processed and distributed by percentage
        data = self.sections_distribution_by_percentage(data)

        # Define the color map for the categories
        color_map = {
            "não conforme": "#E41A1C",  # Red
            "parcial conforme": "#377EB8",  # Blue
            "conforme": "#4DAF4A",  # Green
            "supera": "#FF7F00"  # Orange
        }

        # Convert data to DataFrame and transpose
        df = pd.DataFrame(data).T

        desired_cols = ["supera", "conforme", "parcial conforme", "não conforme"]
        existing_cols = [col for col in desired_cols if col in df.columns]
        df = df[existing_cols]

        df = df.reset_index()
        # Melt the DataFrame to long format for plotting
        df_melted = df.melt(
            id_vars=["index"],
            value_vars=existing_cols,
            var_name="Categoria",
            value_name="Quantidade",
        )

        # Rename the columns for clarity
        df_melted = df_melted.rename(columns={"index": "Setor"})

        # Set up the plot
        plt.figure(figsize=(18, 6))

        # Create the grouped bar plot and capture the axis object
        ax = sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=df_melted,
            palette=color_map,  # Use the custom color map
            edgecolor="black",
        )

        # Annotate bars with percentages
        for p in ax.patches:
            height = p.get_height()
            if height == 0:
                continue
            ax.annotate(
                f"{p.get_height():.1f}%",
                (p.get_x() + p.get_width() / 2.0, p.get_height()),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )

        # Customize the plot
        plt.grid(True, axis="y", linestyle="--", alpha=0.7)
        plt.title(title, fontsize=16)
        plt.ylabel("Distribuição de Respostas")
        plt.xlabel("Seção")

        # Create custom legend handles to match the color map
        handles = []
        for status in existing_cols:
            color = color_map[status]
            handles.append(
                plt.Line2D(
                    [0],
                    [0],
                    marker="o",
                    color="w",
                    markerfacecolor=color,
                    markersize=10,
                    label=status.capitalize(),
                )
            )

        # Add the custom legend to the plot
        plt.legend(handles=handles, title="Categorias", loc="upper right")

        # Adjust layout for better spacing
        plt.tight_layout()

        # Save the figure to a buffer
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)
        return buf


    def plot_grouped_bar_split_section_1(self, data, title):
        # Define the fixed color map for the categories
        color_map = {
            "não conforme": "#E41A1C",  # Red
            "parcial conforme": "#377EB8",  # Blue
            "conforme": "#4DAF4A",      # Green
            "supera": "#FF7F00"         # Orange
        }

        # Convert the dictionary into a DataFrame (transpose to have sections as rows)
        df = pd.DataFrame(data).T
        desired_cols = ["supera", "conforme", "parcial conforme", "não conforme"]
        existing_cols = [col for col in desired_cols if col in df.columns]

        # Reset index so that the index becomes a column; later, we rename it to "Setor"
        df = df.reset_index()

        # Melt the DataFrame for Seaborn
        df_melted = df.melt(
            id_vars=["index"],
            value_vars=existing_cols,
            var_name="Categoria",
            value_name="Quantidade",
        )
        # Rename 'index' to 'Setor'
        df_melted = df_melted.rename(columns={"index": "Setor"})

        # Split the dataset into two halves based on unique "Setor" values
        unique_setors = df_melted["Setor"].unique()
        mid_point = len(unique_setors) // 2
        first_half_setors = unique_setors[:mid_point]
        second_half_setors = unique_setors[mid_point:]

        first_half = df_melted[df_melted["Setor"].isin(first_half_setors)]
        second_half = df_melted[df_melted["Setor"].isin(second_half_setors)]

        # Precompute the total "Quantidade" for each Setor in both halves
        group_sums_first = first_half.groupby("Setor")["Quantidade"].sum().to_dict()
        group_sums_second = second_half.groupby("Setor")["Quantidade"].sum().to_dict()

        # Set up the figure with two subplots
        fig, axes = plt.subplots(2, 1, figsize=(15, 10), sharey=True)

        # --- FIRST SUBPLOT ---
        # Enforce order for x (Setor) and hue (Categoria) to match our DataFrame order
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=first_half,
            palette=color_map,
            edgecolor="black",
            ax=axes[0],
            order=first_half["Setor"].unique(),  # fixed x order
            hue_order=existing_cols,             # fixed hue order
        )
        axes[0].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[0].set_title(title + " (Parte 1)", fontsize=16)
        axes[0].set_ylabel("Quantidade de Respostas")
        axes[0].set_xlabel("Seção")

        # Annotate each bar relative to its group total (skip if value is zero)
        for i, patch in enumerate(axes[0].patches):
            height = patch.get_height()
            if height == 0:
                continue
            # The order of patches corresponds to the order of rows in first_half
            row = first_half.iloc[i]
            setor = row["Setor"]
            total = group_sums_first[setor]
            percent = (height / total) * 100 if total else 0
            axes[0].annotate(
                f"{percent:.1f}%",
                (patch.get_x() + patch.get_width() / 2.0, height),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )

        # --- SECOND SUBPLOT ---
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=second_half,
            palette=color_map,
            edgecolor="black",
            ax=axes[1],
            order=second_half["Setor"].unique(),
            hue_order=existing_cols,
        )
        axes[1].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[1].set_title(title + " (Parte 2)", fontsize=16)
        axes[1].set_ylabel("Quantidade de Respostas")
        axes[1].set_xlabel("Seção")

        for i, patch in enumerate(axes[1].patches):
            height = patch.get_height()
            if height == 0:
                continue
            row = second_half.iloc[i]
            setor = row["Setor"]
            total = group_sums_second[setor]
            percent = (height / total) * 100 if total else 0
            axes[1].annotate(
                f"{percent:.1f}%",
                (patch.get_x() + patch.get_width() / 2.0, height),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )   

        # Create custom legend handles using the fixed color map
        handles = [
            plt.Line2D(
                [0],
                [0],
                marker="o",
                color="w",
                markerfacecolor=color_map[status],
                markersize=10,
                label=status.capitalize(),
            )
            for status in ["supera", "conforme", "parcial conforme", "não conforme"]
        ]

        # Add the custom legend to both subplots
        axes[0].legend(handles=handles, title="Legenda", loc="upper right")
        axes[1].legend(handles=handles, title="Legenda", loc="upper right")

        # Adjust layout for better spacing
        plt.tight_layout()

        # Save the figure to a buffer
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)

        return buf



    # Rename 'index' column to 'Set

    def plot_grouped_bar_split_section_2(self, data, title):
        # Define the fixed color map for the categories
        color_map = {
            "não conforme": "#E41A1C",   # Red
            "parcial conforme": "#377EB8",  # Blue
            "conforme": "#4DAF4A",      # Green
            "supera": "#FF7F00"         # Orange
        }

        # Convert the dictionary into a DataFrame (transpose to have sections as rows)
        df = pd.DataFrame(data).T
        desired_cols = ["supera", "conforme", "parcial conforme", "não conforme"]
        existing_cols = [col for col in desired_cols if col in df.columns]

        # Reorder columns (just to keep them in a known order) and reset index
        df = df[existing_cols]
        df = df.reset_index()  # 'index' will become the "Setor" column

        # Melt the DataFrame for Seaborn
        df_melted = df.melt(
            id_vars=["index"],
            value_vars=existing_cols,
            var_name="Categoria",
            value_name="Quantidade",
        )

        # Rename 'index' column to 'Setor'
        df_melted = df_melted.rename(columns={"index": "Setor"})

        # Split the dataset into three parts (roughly equal, or as needed)
        unique_setors = df_melted["Setor"].unique()
        first_point = len(unique_setors) // 3
        second_point = first_point * 2

        # Make sure slicing is correct even if there's a remainder
        first_part_setors = unique_setors[:first_point]
        second_part_setors = unique_setors[first_point:second_point]
        third_part_setors = unique_setors[second_point:]

        first_part = df_melted[df_melted["Setor"].isin(first_part_setors)]
        second_part = df_melted[df_melted["Setor"].isin(second_part_setors)]
        third_part = df_melted[df_melted["Setor"].isin(third_part_setors)]

        # Precompute sums for each part by "Setor" so we know how to calculate percentages
        group_sums_first = first_part.groupby("Setor")["Quantidade"].sum().to_dict()
        group_sums_second = second_part.groupby("Setor")["Quantidade"].sum().to_dict()
        group_sums_third = third_part.groupby("Setor")["Quantidade"].sum().to_dict()

        # Set up the figure with three subplots
        fig, axes = plt.subplots(3, 1, figsize=(20, 10), sharey=True)

        # --- FIRST SUBPLOT ---
        # We fix the order of x (Setor) and hue (Categoria) so that
        # the order of patches matches the order of rows in first_part.
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=first_part,
            palette=color_map,
            edgecolor="black",
            ax=axes[0],
            order=first_part["Setor"].unique(),  # ensures consistent x ordering
            hue_order=existing_cols,            # ensures consistent hue ordering
        )
        axes[0].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[0].set_title(title + " (Parte 1)", fontsize=16)
        axes[0].set_ylabel("Quantidade de Respostas")
        axes[0].set_xlabel("Seção")

        # Annotate bars for the first subplot (relative to each Setor)
        # We rely on the fact that seaborn plots in the same order we specified.
        for i, patch in enumerate(axes[0].patches):
            height = patch.get_height()
            if height == 0:
                continue

            # Find the corresponding row in first_part
            row = first_part.iloc[i]
            setor = row["Setor"]
            total_for_setor = group_sums_first[setor]
            percent = (height / total_for_setor) * 100

            axes[0].annotate(
                f"{percent:.1f}%",
                (patch.get_x() + patch.get_width() / 2.0, height),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )

        # --- SECOND SUBPLOT ---
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=second_part,
            palette=color_map,
            edgecolor="black",
            ax=axes[1],
            order=second_part["Setor"].unique(),
            hue_order=existing_cols,
        )
        axes[1].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[1].set_title(title + " (Parte 2)", fontsize=16)
        axes[1].set_ylabel("Quantidade de Respostas")
        axes[1].set_xlabel("Seção")

        for i, patch in enumerate(axes[1].patches):
            height = patch.get_height()
            if height == 0:
                continue

            row = second_part.iloc[i]
            setor = row["Setor"]
            total_for_setor = group_sums_second[setor]
            percent = (height / total_for_setor) * 100

            axes[1].annotate(
                f"{percent:.1f}%",
                (patch.get_x() + patch.get_width() / 2.0, height),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )

        # --- THIRD SUBPLOT ---
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=third_part,
            palette=color_map,
            edgecolor="black",
            ax=axes[2],
            order=third_part["Setor"].unique(),
            hue_order=existing_cols,
        )
        axes[2].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[2].set_title(title + " (Parte 3)", fontsize=16)
        axes[2].set_ylabel("Quantidade de Respostas")
        axes[2].set_xlabel("Seção")

        for i, patch in enumerate(axes[2].patches):
            height = patch.get_height()
            if height == 0:
                continue

            row = third_part.iloc[i]
            setor = row["Setor"]
            total_for_setor = group_sums_third[setor]
            percent = (height / total_for_setor) * 100

            axes[2].annotate(
                f"{percent:.1f}%",
                (patch.get_x() + patch.get_width() / 2.0, height),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )

        # Create custom legend handles using the fixed color map
        handles = [
            plt.Line2D(
                [0],
                [0],
                marker="o",
                color="w",
                markerfacecolor=color_map[status],
                markersize=10,
                label=status.capitalize(),
            )
            for status in existing_cols
        ]

        # Add the custom legend to all subplots
        axes[0].legend(handles=handles, title="Legenda", loc="upper right")
        axes[1].legend(handles=handles, title="Legenda", loc="upper right")
        axes[2].legend(handles=handles, title="Legenda", loc="upper right")

        # Adjust layout for better spacing
        plt.tight_layout()

        # Save the figure to a buffer
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)

        return buf



    def plot_line_graph_from_queryset(self, queryset, title):
        # Define the fixed color map for the categories
        color_map = {
            "não conforme": "#E41A1C",  # Red
            "parcial conforme": "#377EB8",  # Blue
            "conforme": "#4DAF4A",  # Green
            "supera": "#FF7F00"  # Orange
        }

        # Prepare data containers
        dates = []
        conforme = []
        nao_conforme = []
        parcial_conforme = []
        supera = []

        # Collect data from the queryset
        for record in queryset:
            dates.append(record.date)
            distribution = self.awnser_distribution_by_percentage(record.ona_answer_total_distribution)
            conforme.append(distribution.get('conforme', 0))
            nao_conforme.append(distribution.get('não conforme', 0))
            parcial_conforme.append(distribution.get('parcial conforme', 0))
            supera.append(distribution.get('supera', 0))

        # Create a DataFrame from the data
        df = pd.DataFrame({
            "Date": dates,
            "Conforme": conforme,
            "Não conforme": nao_conforme,
            "Parcial conforme": parcial_conforme,
            "Supera": supera,
        })

        # Set up the figure
        plt.figure(figsize=(12, 6))

        # Plot each category as a line with dots at the data points using the color_map
        plt.plot(df['Date'], df['Conforme'], label='Conforme', marker='o', linestyle='-', color=color_map["conforme"])
        plt.plot(df['Date'], df['Não conforme'], label='Não conforme', marker='o', linestyle='-', color=color_map["não conforme"])
        plt.plot(df['Date'], df['Parcial conforme'], label='Parcial conforme', marker='o', linestyle='-', color=color_map["parcial conforme"])
        plt.plot(df['Date'], df['Supera'], label='Supera', marker='o', linestyle='-', color=color_map["supera"])

        # Customize the plot
        plt.title(title, fontsize=16)
        plt.xlabel('Data', fontsize=14)
        plt.ylabel('Quantidade de Respostas', fontsize=14)

        # Set the x-axis ticks to only show the dates where data exists
        plt.gca().xaxis.set_major_locator(mdates.WeekdayLocator())  # Optional: to show ticks for each week
        plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%d/%m/%y'))

        # Set the x-ticks to only display the dates in the dataset
        plt.xticks(df['Date'], rotation=45)

        # Show grid lines
        plt.grid(True)
        desired_cols = ["supera", "conforme", "parcial conforme", "não conforme"]
        existing_cols = [col for col in desired_cols if col in df.columns]
        # Create custom legend handles using the fixed color map
        handles = [
            plt.Line2D([0], [0], marker='o', color='w', markerfacecolor=color_map[status], markersize=10, label=status)
            for status in  existing_cols
        ]

        # Add the custom legend to the plot
        plt.legend(handles=handles, title='Categorias', loc='upper right')

        # Show the plot
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)
        return buf
    
    def awnser_distribution_by_percentage(self, answer_distribution: dict) -> dict:

        if "não aplicável" in answer_distribution.keys():
            del answer_distribution["não aplicável"]


        print(answer_distribution)
        total = sum(answer_distribution.values())
        for key, value in answer_distribution.items():
            answer_distribution[key] = round((value / total) * 100, 2)

        return answer_distribution

    def sections_distribution_by_percentage(self, sections_distribution: dict) -> dict:
        # For each section, apply the answer distribution function

        for key, value in sections_distribution.items():
            sections_distribution[key] = self.awnser_distribution_by_percentage(value)

        return sections_distribution

class PDFReportGenerator:
    def __init__(self, filename):
        self.filename = filename

        self.save_path = f"/tmp/relatorio_formulario_{timezone.now()}.pdf"
        self.doc = SimpleDocTemplate(
            self.save_path,
            pagesize=letter
        )
        self.story = []  # This will hold the content of the PDF
        self.metrics = MetricsCalculator()
        self.graphs = GraphsGenerator()
        
      
        
    
    def create_pdf_report_for_subsection(self, evaluator_name: str, answers: ONAFormAnswered, evaluator_id: int):
        
    # Build the PDF content
        sections = answers.ona_form.ONA_sections.all()
        section_title = sections[0].section_title    
        subsections = answers.answered_sections.all()
        subsections = subsections[0].answered_subsections.all()
        subsection_title = subsections[0].form_subsection.subsection_title    
        
        self.add_title(
            evaluator_name=evaluator_name,
            section_title=section_title,
            subsection_title=subsection_title,
        )
        metrics = self.metrics.get_ona_form_average_distribution(ona_form=answers)
        
        form_distribution_img = self.graphs.plot_bar_plot(
            data=metrics["ONA answer Distribution"],
            title="Distribuição das Respostas no formulario",
        )
        self.insert_image_center(image=form_distribution_img)
        self.display_answers_with_comments(questions_answers=metrics["Answers with comments"])
        
        # Generate a timestamped save path
       
      
        
        
        # Reinitialize self.doc with the save path so that build() writes directly to this file
     
        self.doc.build(self.story)
        evaluator = Evaluator.objects.filter(
            id=evaluator_id
        ).first()
        evaluator_email=evaluator.email
        # Send the email with the generated PDF attached
        self.send_email_report(report_path=self.save_path, evaluator_email = evaluator_email)

    def send_email_report(self, report_path, evaluator_email):
        subject = "Generated Report"
        body = "Segue em anexo o relatório completo da avaliação."
        email_list = [
            email.strip() 
            for email in os.getenv("EMAIL_LIST", "").split(",") 
            if email.strip()
        ]
        
        email = EmailMessage(
            subject=subject,
            body=body,
            from_email=os.getenv("EMAIL_HOST_USER"), #+ evaluator_email,
            to=email_list,
        )
       
# Save your file to temp_file_path and attach it to your email

        with open(report_path, "rb") as file:
            email.attach_file(report_path)
        
        try:
            email.send()
            print(f"Email sent to {email_list}")
        except Exception as e:
            print(f"Failed to send email: {str(e)}")
            raise
            
    def add_title(self, evaluator_name: str, section_title: str, subsection_title: str) -> None:
        # Register the DejaVuSans font (supports extended Latin characters)
        pdfmetrics.registerFont(TTFont('DejaVuSans', 'report/helpers/DejaVuSans-Bold.ttf'))
        
        # Get default styles and set their font to DejaVuSans
        styles = getSampleStyleSheet()
        styles['Normal'].fontName = 'DejaVuSans'
        styles['Title'].fontName = 'DejaVuSans'
        
        # Build the title text
        title = (
            f"Relatório de preenchimento da subseção {subsection_title} "
            f"na seção {section_title}, feito por {evaluator_name}"
        )
        
        # Create a Paragraph using the updated 'Title' style
        title_paragraph = Paragraph(title, styles['Title'])
        
        # Add the title paragraph to your story
        self.story.append(title_paragraph)

    def display_answers_with_comments(self, questions_answers):
        self.story.append(PageBreak())
        styles = getSampleStyleSheet()
        
        
        questions_answers = sorted(questions_answers, key=lambda x: x.question.question_id)
        for question in questions_answers:
            data = []

            # Add header row (with column names)
            header = [question.question.question_id, 'Core', "Resposta", "Comentario"]
            data.append(header)
            page_width, page_height = letter
            left_margin, right_margin = 4 , 4   
            table_width = page_width - left_margin - right_margin
        
            question_desc = f"{question.question.description}"
            core = "Sim" if question.question.core else "Não"
            core_info = f"{core}"
            answer_info = f"{question.answer}"
            comment_info = f"{question.comment}"

            
            row = [
                Paragraph(question_desc, styles['Normal']),
                Paragraph(core_info, styles['Normal']),
                Paragraph(answer_info, styles['Normal']),
                Paragraph(comment_info, styles['Normal'])
            ]
            data.append(row)
            col_widths = [table_width * 0.20, table_width * 0.20, table_width * 0.20, table_width * 0.20]
        
            table = Table(data, colWidths=col_widths, spaceAfter=20) 
            table.setStyle(TableStyle([
                ('GRID', (0, 0), (-1, -1), 1, colors.black),  
                ('BACKGROUND', (0, 0), (-1, 0), colors.blue), 
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white), 
                ('ALIGN',(0,-1),(-1,-1),'CENTER'),
                ('VALIGN',(0,-1),(-1,-1),'MIDDLE'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'), 
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12), 
                ('TOPPADDING', (0, 1), (-1, -1), 8),  
                ('LEFTPADDING', (0, 0), (-1, -1), 5), 
                ('RIGHTPADDING', (0, 0), (-1, -1), 5),
                # ('SPACEAFTER', (0, 0), (-1, -1), 50)
            ]))

            # Add the table to the story (content of the PDF)
            self.story.append(table)
            self.story.append(Paragraph("<br />", styles['Normal']))
    def insert_image_center(self, image: BytesIO) -> None:
        # The width and height for the image in the PDF
        img_width, img_height = 400, 400
     
        
        # Use the Image flowable to insert the image into the story
        img = Image(image, width=img_width, height=img_height)
        img.hAlign = 'CENTER'  # This centers the image in the PDF
        
        # Add the image flowable to the story
        self.story.append(img)

class PowerPointReportGenerator:
    def __init__(self):
        self.template_ona_path = "report/helpers/template-ona-report.pptx"
        self.section_indexes = {
            "1 SEÇÃO - GESTÃO ORGANIZACIONAL": 1,
            "2 SEÇÃO - ATENÇÃO AO PACIENTE": 4,
            "3 SEÇÃO - DIAGNÓSTICO E TERAPÊUTICA": 7,
            "4 SEÇÃO - GESTÃO DE APOIO": 10,
        }
        self.presentation = Presentation(self.template_ona_path)
        self.graphs = GraphsGenerator()

    def insert_plot_image_in_slide(self, image_buffer, slide_index, image_type):
        left, top, width, height = self.generate_images_coords(image_type)
        slide = self.presentation.slides[slide_index]
        slide.shapes.add_picture(image_buffer, left, top, width, height)
        return self.presentation

    def generate_images_coords(self, image_type):
        
        left = Inches(1.5)
        top = Inches(1.3)
        width = Inches(10)
        height = Inches(6)
        return left, top, width, height

    def add_section_images(self, sections_data):
        for section, resuls in sections_data.items():
            bar_plot_img = self.graphs.plot_bar_plot(resuls, section)
            section = section.strip()
            
            slide_index = self.section_indexes[section]
 
            prs = self.insert_plot_image_in_slide(bar_plot_img, slide_index, "section")
        return prs

    def add_subsection_images(self, subsection_data):
        for subsection, results in subsection_data.items():
            slide_index = self.section_indexes[subsection] + 2
            if slide_index == 3:
                grouped_bar_img = self.graphs.plot_grouped_bar_split_section_1(
                    results, subsection
                )
            elif slide_index == 6:
                grouped_bar_img = self.graphs.plot_grouped_bar_split_section_2(
                    results, subsection
                )
            else:
                grouped_bar_img = self.graphs.plot_grouped_bar(results, subsection)
            prs = self.insert_plot_image_in_slide(
                grouped_bar_img, slide_index, "subsection"
            )
        return prs

    def add_plot_line_image(self, hospital):
        queryset = ONAFormDistribution.objects.filter(
            hospital=hospital
        )
        self.graphs.plot_line_graph_from_queryset(
            queryset=queryset,
            title="Evolução das Respostas"
        )
        plot_line_img = self.graphs.plot_line_graph_from_queryset(
            queryset=queryset,
            title="Evolução das Respostas"
        )
        prs = self.insert_plot_image_in_slide(
            image_buffer=plot_line_img,
            slide_index=13,
            image_type="plot-line"
        )
        return prs
    
    def add_core_distributin_image(self, core_ditribution):

        for section, resuls in core_ditribution.items():
            
            bar_plot_img = self.graphs.plot_bar_plot(resuls, section)
            section = section.strip()
            
            slide_index = self.section_indexes[section] +1
 
            prs = self.insert_plot_image_in_slide(bar_plot_img, slide_index, "section")
        return prs




        

    def make_report(self, data, report_name, hospital):
        subsections = data["Subsections Distribution"]
        sections = data["Sections Distribution"]
        core_distribution = data['Core Questions Distribution']
   
        self.add_section_images(sections)
        self.add_core_distributin_image(core_distribution)
        self.add_subsection_images(subsections)
        self.add_plot_line_image(hospital)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        path = f"/tmp/apresentacao_do_dia_{timestamp}.pptx"
        self.presentation.save(path)
        self.send_email(path)

    def send_email(self, report_path):
        subject = "Generated Report"
        body = "Segue em anexo o relatório completo da avaliação."
        email_list = [
            email.strip() 
            for email in os.getenv("EMAIL_LIST", "").split(",") 
            if email.strip()
        ]
        
        # Create the email message
        email = EmailMessage(
            subject=subject,
            body=body,
            from_email=os.getenv("EMAIL_HOST_USER"),  # Use the sender email from settings
            to=email_list,
        )

        # Attach the report file
        with open(report_path, "rb") as file:
            email.attach_file(report_path)

        # Send the email
        try:
            email.send()  # This sends the email using Django's configured email backend
            print(f"Email sent to {email_list}")
        except Exception as e:
            print(f"Error sending email: {e}")


class MetricsCalculator:
    def __get_subsection_questions(
        self, subsection: FormSubsectionAnswered
    ) -> list[QuestionAnswer]:
        leve1_questions = subsection.answered_questions_level_1.all()
        level2_questions = subsection.answered_questions_level_2.all()

        questions_answers = leve1_questions.union(level2_questions)
        return questions_answers

    def __get_questions_average_distribution(
        self, questions_list: list[QuestionAnswer]
    ) -> dict[str, int]:
        questions_answers = [question.answer for question in questions_list]
        answers_distribution = Counter(questions_answers)
        return answers_distribution

    def __get_subsection_average_distribution(
        self, subsection_questions: list[QuestionAnswer]
    ) -> dict[str, dict[str, int]]:
        subsection_distribution = self.__get_questions_average_distribution(
            questions_list=subsection_questions
        )
        return dict(subsection_distribution)
    
    def __get_core_questions(self, questions_list: list[QuestionAnswer]) -> list[QuestionAnswer]:
        core_questions = []
        for question in questions_list:
            if question.question.core: 
                core_questions.append(question)
        return core_questions 
                
    def __get_section_and_subsection_answers_distribution(
        self, section: FormSectionAnswered
    ) -> dict[str, int]:
        subsections = section.answered_subsections.all()
   
        level3_questions = section.answered_questions_level_3.all()

        subsection_answers = level3_questions
        
        subsection_distribution = {}
        for subsction in subsections:
            subsection_title = subsction.form_subsection.subsection_title
            subsection_questions_level_1_and_level_2 = self.__get_subsection_questions(
                subsction
            )
            subsection_distribution[subsection_title] = (
                self.__get_subsection_average_distribution(
                    subsection_questions_level_1_and_level_2
                )
            )

            subsection_answers = subsection_answers.union(
                subsection_questions_level_1_and_level_2
            )

        core_questions = self.__get_core_questions(subsection_answers)
        
        section_distribution = self.__get_questions_average_distribution(
            subsection_answers
        )
        core_questions_distribution = self.__get_questions_average_distribution(
            core_questions
        )
       
        



        return section_distribution, subsection_distribution, subsection_answers, core_questions_distribution

    def __get_ona_form_total_metrics(
        self, distribution_by_section: dict
    ) -> dict[str, int]:
        total_counts = Counter()
        for section, counts in distribution_by_section.items():
            total_counts.update(counts)
        return dict(total_counts)

    def get_ona_form_average_distribution(
        self, ona_form: ONAFormAnswered
    ) -> dict[str, dict[str, int]]:
        sections = ona_form.answered_sections.all()

        distribution_by_section = {}
        distribution_by_subsections = {}
        distribution_by_core = {}
        for section in sections:
            section_name = section.form_section.section_title
            section_distribution, subsection_distribution, section_answers_with_comments, core_questions_distribution = (
                self.__get_section_and_subsection_answers_distribution(section)
            )

            distribution_by_section[section_name] = dict(section_distribution)

            distribution_by_subsections[section_name] = dict(subsection_distribution)
            distribution_by_core[section_name] = dict(core_questions_distribution)

        total_distribution = self.__get_ona_form_total_metrics(distribution_by_section)

        metrics = {
            "Subsections Distribution": distribution_by_subsections,
            "Sections Distribution": distribution_by_section,
            "ONA answer Distribution": total_distribution,
            "Answers with comments" : section_answers_with_comments,
            "Core Questions Distribution" : distribution_by_core
        }
 

       
    
        return metrics
    

    def create_unified_form_metrics(self, ona_form_queryset: QuerySet):
        combined_distribution = {
            'não conforme': 0,
            'parcial conforme': 0,
            'conforme': 0,
            'supera': 0
        }
        combined_sections_distribution = {}
        combined_subsections_distribution = {}
        combined_core_distribution = {}
        for ona_form in ona_form_queryset:
        
            metrics = self.get_ona_form_average_distribution(ona_form)
            combined_distribution = self.update_combined_distribution(
                ona_answer_distribution=metrics['ONA answer Distribution'],
                combined_distribution=combined_distribution
            )
            combined_sections_distribution = self.update_combine_sections_distribution(
                sections_distribution=metrics["Sections Distribution"],
                combined_sections_distribution=combined_sections_distribution,
            )
            combined_subsections_distribution = self.update_combine_subsections_distribution(
                subsections_distribution_by_sections = metrics["Subsections Distribution"],
                combined_sections_distribution = combined_subsections_distribution
            )
       
            combined_core_distribution = self.update_combine_sections_distribution(
                sections_distribution = metrics['Core Questions Distribution'],
                combined_sections_distribution = combined_core_distribution
            )

         
        combined_metrics = {
            "Subsections Distribution": dict(sorted(combined_subsections_distribution.items())),
            "Sections Distribution": dict(sorted(combined_sections_distribution.items())),
            "ONA answer Distribution": dict(sorted(combined_distribution.items())),
            "Core Questions Distribution": dict(sorted(combined_core_distribution.items()))
        }
        return combined_metrics 
            
    def update_combined_distribution(self, ona_answer_distribution: dict, combined_distribution:dict) -> dict:
        for key, value in ona_answer_distribution.items():
            if key in combined_distribution:
                combined_distribution[key] += value
            else:
                combined_distribution[key] = value
        return combined_distribution
    
    def update_combine_sections_distribution(self,  sections_distribution:dict, combined_sections_distribution:dict):

        for section_name, distribution in sections_distribution.items():
            if section_name in combined_sections_distribution:
                # Use Counter to add the counts from both dictionaries
                combined_sections_distribution[section_name] = dict(
                    Counter(distribution) + Counter(combined_sections_distribution[section_name])
                )
            else:
                # If the section doesn't exist, just add it.
                combined_sections_distribution[section_name] = distribution.copy()  # copy to avoid mutating the input

        return combined_sections_distribution
    
    # def update_combined_core_distribution(self, ona_answer_distribution: dict, combined_distribution: dict) -> dict:
    #     for section, distribution in ona_answer_distribution.items():
    #         if section in combined_distribution:
    #             ona_answer_distribution[section] =   Counter(distribution) + Counter(ona_answer_distribution[section])
               
    #         else:
    #             # Use a copy of counts to avoid aliasing issues if the caller reuses the dictionary
    #             combined_distribution[section] = distribution.copy()
    #     return combined_distribution

    
    def update_combine_subsections_distribution(self, subsections_distribution_by_sections: dict, combined_sections_distribution: dict):
        for section_name, subsection_distribution in subsections_distribution_by_sections.items():
            if section_name in combined_sections_distribution:
                # If the section exists in the combined dictionary, update its subsections
                for subsection_name, distribution in subsection_distribution.items():
                    if subsection_name in combined_sections_distribution[section_name]:
                        # Use Counter to sum the values of matching subsections
                        combined_sections_distribution[section_name][subsection_name] = dict(
                            Counter(combined_sections_distribution[section_name][subsection_name]) + Counter(distribution)
                        )
                    else:
                        # If the subsection does not exist in the combined section, add it
                        combined_sections_distribution[section_name][subsection_name] = distribution
            else:
                # If the section does not exist in the combined dictionary, add it directly
                combined_sections_distribution[section_name] = subsection_distribution

        return combined_sections_distribution





