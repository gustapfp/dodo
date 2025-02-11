
from collections import Counter
import matplotlib.pyplot as plt
from datetime import datetime
import seaborn as sns
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from django.core.mail import EmailMessage
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

class GraphsGenerator:
    def plot_bar_plot(self, data, title):
        data = self.awnser_distribution_by_percentage(data)

        colors = ["#E41A1C", "#377EB8", "#4DAF4A", "#FF7F00"]

        df = pd.DataFrame(list(data.items()), columns=["Status", "Percentage"])
        df["Status"] = df["Status"].str.capitalize()

        plt.figure(figsize=(10, 6))

        ax = sns.barplot(
            x="Status", y="Percentage", data=df, palette=colors, edgecolor="black"
        )

        plt.ylabel("Conformidade (%)")
        plt.xlabel("Respostas", labelpad=15)
        plt.title(title)

        plt.grid(True, axis="y", linestyle="--", alpha=0.7)

        plt.ylim(0, 100)

        for p in ax.patches:
            ax.annotate(
                f"{p.get_height():.1f}%",
                (p.get_x() + p.get_width() / 2.0, p.get_height()),
                ha="center",
                va="center",
                fontsize=12,
                color="black",
                fontweight="bold",
                xytext=(0, 10),
                textcoords="offset points",
            )

        colors = ax.patches

        handles = []
        categories = df["Status"].tolist()
        for color, status in zip(colors, categories):
            handles.append(
                plt.Line2D(
                    [0],
                    [0],
                    marker="o",
                    color="w",
                    markerfacecolor=color.get_facecolor(),
                    markersize=10,
                    label=status,
                )
            )

        plt.legend(handles=handles, title="Legenda")
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)
        return buf

    def plot_grouped_bar(self, data, title):
        data = self.sections_distribution_by_percentage(data)
        colors = [
            "#E41A1C",
            "#377EB8",
            "#4DAF4A",
            "#FF7F00",
        ]  # Red, Blue, Green, Orange

        df = pd.DataFrame(data).T
        df = df[["conforme", "não conforme", "supera", "parcial conforme"]]
        df = df.reset_index()

        df_melted = df.melt(
            id_vars=["index"],
            value_vars=["supera", "conforme", "parcial conforme", "não conforme"],
            var_name="Categoria",
            value_name="Quantidade",
        )

        df_melted = df_melted.rename(columns={"index": "Setor"})

        plt.figure(figsize=(18, 6))

        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=df_melted,
            palette=colors,
            edgecolor="black",
        )
        plt.grid(True, axis="y", linestyle="--", alpha=0.7)

        # Customize the plot
        plt.title(title, fontsize=16)
        plt.ylabel("Quantidade de Respostas")
        plt.xlabel("Seção")

        # Create custom legend handles using plt.Line2D with hard-coded colors
        handles = []
        categories = ["supera", "conforme", "parcial conforme", "não conforme"]

        # Iterate through the bars and create a custom legend entry for each
        for i, status in enumerate(categories):
            # Use the corresponding color from the hard-coded colors list
            color = colors[i]
            # Create the legend entry using a square marker ('o')
            handles.append(
                plt.Line2D(
                    [0],
                    [0],
                    marker="o",
                    color="w",
                    markerfacecolor=color,
                    markersize=10,
                    label=status.capitalize(),
                )
            )

        # Add the custom legend
        plt.legend(handles=handles, title="Categorias", loc="upper right")

        # Adjust layout for better spacing
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)
        return buf

    def plot_grouped_bar_split_section_1(self, data, title):
        # Define the hard-coded colors from the Set1 palette
        colors = [
            "#E41A1C",
            "#377EB8",
            "#4DAF4A",
            "#FF7F00",
        ]  # Set1 colors: red, blue, green, orange

        # Convert the dictionary into a DataFrame
        df = pd.DataFrame(data).T  # Transpose to have sections as rows
        df = df[
            ["conforme", "não conforme", "supera", "parcial conforme"]
        ]  # Ensure correct column order
        df = df.reset_index()  # Reset the index to have a 'Setor' column for seaborn

        # Melt the dataframe for seaborn
        df_melted = df.melt(
            id_vars=["index"],
            value_vars=["supera", "conforme", "parcial conforme", "não conforme"],
            var_name="Categoria",
            value_name="Quantidade",
        )

        # Rename 'index' column to 'Setor'
        df_melted = df_melted.rename(columns={"index": "Setor"})

        # Split the dataset into two halves
        mid_point = len(df_melted["Setor"].unique()) // 2
        first_half = df_melted[df_melted["Setor"].isin(df["index"][:mid_point])]
        second_half = df_melted[df_melted["Setor"].isin(df["index"][mid_point:])]

        # Set up the figure with two subplots
        fig, axes = plt.subplots(2, 1, figsize=(15, 10), sharey=True)

        # First subplot
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=first_half,
            palette=colors,
            edgecolor="black",
            ax=axes[0],
        )
        axes[0].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[0].set_title(title + " (Parte 1)", fontsize=16)
        axes[0].set_ylabel("Quantidade de Respostas")
        axes[0].set_xlabel("Seção")

        # Second subplot
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=second_half,
            palette=colors,
            edgecolor="black",
            ax=axes[1],
        )
        axes[1].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[1].set_title(title + " (Parte 2)", fontsize=16)
        axes[1].set_ylabel("Quantidade de Respostas")
        axes[1].set_xlabel("Seção")

        # Create custom legend handles
        handles = [
            plt.Line2D(
                [0],
                [0],
                marker="o",
                color="w",
                markerfacecolor=colors[i],
                markersize=10,
                label=status.capitalize(),
            )
            for i, status in enumerate(
                ["supera", "conforme", "parcial conforme", "não conforme"]
            )
        ]

        # Add the custom legend to the first plot only
        axes[0].legend(handles=handles, title="Legenda", loc="upper right")
        axes[1].legend(handles=handles, title="Legenda", loc="upper right")

        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)
        return buf

    def plot_grouped_bar_split_section_2(self, data, title):
        # Define the hard-coded colors from the Set1 palette
        colors = [
            "#E41A1C",
            "#377EB8",
            "#4DAF4A",
            "#FF7F00",
        ]  # Set1 colors: red, blue, green, orange

        # Convert the dictionary into a DataFrame
        df = pd.DataFrame(data).T  # Transpose to have sections as rows
        df = df[
            ["conforme", "não conforme", "supera", "parcial conforme"]
        ]  # Ensure correct column order
        df = df.reset_index()  # Reset the index to have a 'Setor' column for seaborn

        # Melt the dataframe for seaborn
        df_melted = df.melt(
            id_vars=["index"],
            value_vars=["supera", "conforme", "parcial conforme", "não conforme"],
            var_name="Categoria",
            value_name="Quantidade",
        )

        # Rename 'index' column to 'Setor'
        df_melted = df_melted.rename(columns={"index": "Setor"})

        # Split the dataset into two halves
        first_point = len(df_melted["Setor"].unique()) // 3
        second_point = first_point * 2 + 1
        first_part = df_melted[df_melted["Setor"].isin(df["index"][:first_point])]
        second_part = df_melted[
            df_melted["Setor"].isin(df["index"][first_point:second_point])
        ]
        third_part = df_melted[df_melted["Setor"].isin(df["index"][second_point:])]

        # Set up the figure with two subplots
        fig, axes = plt.subplots(3, 1, figsize=(20, 10), sharey=True)

        # First subplot
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=first_part,
            palette=colors,
            edgecolor="black",
            ax=axes[0],
        )
        axes[0].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[0].set_title(title + " (Parte 1)", fontsize=16)
        axes[0].set_ylabel("Quantidade de Respostas")
        axes[0].set_xlabel("Seção")

        # Second subplot
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=second_part,
            palette=colors,
            edgecolor="black",
            ax=axes[1],
        )
        axes[1].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[1].set_title(title + " (Parte 2)", fontsize=16)
        axes[1].set_ylabel("Quantidade de Respostas")
        axes[1].set_xlabel("Seção")

        # Second subplot
        sns.barplot(
            x="Setor",
            y="Quantidade",
            hue="Categoria",
            data=third_part,
            palette=colors,
            edgecolor="black",
            ax=axes[2],
        )
        axes[2].grid(True, axis="y", linestyle="--", alpha=0.7)
        axes[2].set_title(title + " (Parte 2)", fontsize=16)
        axes[2].set_ylabel("Quantidade de Respostas")
        axes[2].set_xlabel("Seção")
        # Create custom legend handles
        handles = [
            plt.Line2D(
                [0],
                [0],
                marker="o",
                color="w",
                markerfacecolor=colors[i],
                markersize=10,
                label=status.capitalize(),
            )
            for i, status in enumerate(
                ["supera", "conforme", "parcial conforme", "não conforme"]
            )
        ]

        # Add the custom legend to the first plot only
        axes[0].legend(handles=handles, title="Legenda", loc="upper right")
        axes[1].legend(handles=handles, title="Legenda", loc="upper right")
        axes[2].legend(handles=handles, title="Legenda", loc="upper right")

        # Adjust layout
        plt.tight_layout()
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format="png")
        plt.close()  # Close the figure to free memory
        buf.seek(0)
        return buf

    def awnser_distribution_by_percentage(self, answer_distribution: dict) -> dict:
        # Remove 'não aplicável' if it exists
        if "não aplicável" in answer_distribution.keys():
            del answer_distribution["não aplicável"]

        # Calculate the percentage distribution
        total = sum(answer_distribution.values())
        for key, value in answer_distribution.items():
            answer_distribution[key] = round((value / total) * 100, 2)

        return answer_distribution

    def sections_distribution_by_percentage(self, sections_distribution: dict) -> dict:
        # For each section, apply the answer distribution function
        for key, value in sections_distribution.items():
            sections_distribution[key] = self.awnser_distribution_by_percentage(value)

        return sections_distribution

class PDFReportGeneator:
    def __init__(self, filename):
        self.width, self.height = letter
        self.pdf_canvas = canvas.Canvas(
            filename=filename,
            pagesize=letter
        )
        self.graphs = GraphsGenerator()
        self.metrics  = MetricsCalculator()
        # self.font = "Helverica"

    def create_pdf_report_for_subsection(self, filename:str, evaluator_name:str, answers):
        self.add_title(
            filename=filename,
            evaluator_name=evaluator_name
        )

        metrics = answers
        # self.metrics.get_ona_form_average_distribution(
        #     ona_form=answers
        # )
    
        
        form_distribution_img = self.graphs.plot_bar_plot(
            data=metrics["ONA answer Distribution"],
            title="Distribuição das Respostas no formulario",
        )

        self.insert_image_center(
            image=form_distribution_img
        )

        self.display_answers_with_comments(
            questions_answers=metrics["Answers with comments"],
            y_position = 300
        )

        self.pdf_canvas.save()


    def add_title(self, filename:str, evaluator_name:str) -> None:
        title = f"Relatório de Preenchimento do Formulário por {evaluator_name}"
        title_width = self.pdf_canvas.stringWidth(title)

        title_x_coord = (self.width - title_width)/2
        title_y_coord = self.height - 100

        self.pdf_canvas.drawString(
            title_x_coord,
            title_y_coord,
            filename
        )

    def display_answers_with_comments(self, questions_answers, y_position):
        for question_with_comment in questions_answers:

            
            #  Description
            question_desc = f"Question Description: {question_with_comment['question']}"
            self.pdf_canvas.drawString(100, y_position, question_desc)
            y_position -= 15

            # Core 
            core_info = f"Core: {question_with_comment['question']}"
            self.pdf_canvas.drawString(100, y_position, core_info)
            y_position -= 15
            
            # Answer to the Question
            question_answer = f"Answer: {question_with_comment['answer']}"
            self.pdf_canvas.drawString(100, y_position, question_answer)
            y_position -= 15
            
            # Comment
            question_comment = f"Comment: {question_with_comment['comment']}"
            self.pdf_canvas.drawString(100, y_position, question_comment)
            y_position -= 30 

       
            if y_position < 60:
                self.pdf_canvas.showPage() 
                y_position = self.height - 40 



    
    def insert_image_center(self, image:BytesIO) -> None:
        img_width, img_height = 200, 200  # You can adjust these dimensions as needed
        img_x_coord = (self.width - img_width) / 2
        img_y_coord = (self.height - img_height) / 2
        
        image = ImageReader(image)
        # Draw the image (adjust the width, height, and position as needed)
        self.pdf_canvas.drawImage(image, img_x_coord, img_y_coord, width=img_width, height=img_height)

class PowerPointReportGenerator:
    def __init__(self):
        self.template_ona_path = "report/helpers/template-ona-report.pptx"
        self.section_indexes = {
            "SEÇÃO 01 - GESTÃO ORGANIZACIONAL": 1,
            "SEÇÃO 02  - ATENÇÃO AO PACIENTE": 3,
            "SEÇÃO 03 - DIAGNÓSTICO E TERAPÊUTICA": 5,
            "SEÇÃO 04 - GESTÃO DE APOIO": 7,
        }
        self.presentation = Presentation(self.template_ona_path)
        self.graphs = GraphsGenerator()

    def insert_plot_image_in_slide(self, image_buffer, slide_index, image_type):
        left, top, width, height = self.generate_images_coords(image_type)
        slide = self.presentation.slides[slide_index]
        slide.shapes.add_picture(image_buffer, left, top, width, height)
        return self.presentation

    def generate_images_coords(self, image_type):
        if image_type == "section":
            left = Inches(5.5)
            top = Inches(1)
            width = Inches(7)
            height = Inches(5)
        else:
            left = Inches(1.5)
            top = Inches(1.3)
            width = Inches(10)
            height = Inches(6)
        return left, top, width, height

    def add_section_images(self, sections_data):
        for section, resuls in sections_data.items():
            bar_plot_img = self.graphs.plot_bar_plot(resuls, section)
            slide_index = self.section_indexes[section]
            prs = self.insert_plot_image_in_slide(bar_plot_img, slide_index, "section")
        return prs

    def add_subsection_images(self, subsection_data):
        for subsection, results in subsection_data.items():
            slide_index = self.section_indexes[subsection] + 1
            if slide_index == 2:
                grouped_bar_img = self.graphs.plot_grouped_bar_split_section_1(
                    results, subsection
                )
            elif slide_index == 4:
                grouped_bar_img = self.graphs.plot_grouped_bar_split_section_2(
                    results, subsection
                )
            else:
                grouped_bar_img = self.graphs.plot_grouped_bar(results, subsection)
            prs = self.insert_plot_image_in_slide(
                grouped_bar_img, slide_index, "subsection"
            )
        return prs

    def make_report(self, data, report_name):
        subsections = data["Subsections Distribution"]
        sections = data["Sections Distribution"]
        self.add_section_images(sections)
        self.add_subsection_images(subsections)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = f"report/presentations_report/{report_name}_{timestamp}.pptx"
        self.presentation.save(path)
        self.send_email(path, "gustavopfpereira30@gmail.com")

    def send_email(self, report_path, recipient_email):
        subject = "Generated Report"
        body = "Please find attached the generated report."

        # Create the email message
        email = EmailMessage(
            subject=subject,
            body=body,
            from_email="gustavopfpereira30@gmail.com",  # Use the sender email from settings
            to=["gustavopfpereira30@gmail.com"],
        )

        # Attach the report file
        with open(report_path, "rb") as file:
            email.attach_file(report_path)

        # Send the email
        try:
            email.send()  # This sends the email using Django's configured email backend
            print(f"Email sent to {recipient_email}")
        except Exception as e:
            print(f"Error sending email: {e}")


class MetricsCalculator:
    def __get_subsection_questions(
        self, subsection
    ):
        leve1_questions = subsection.answered_questions_level_1.all()
        level2_questions = subsection.answered_questions_level_2.all()

        questions_answers = leve1_questions.union(level2_questions)
        return questions_answers

    def __get_questions_average_distribution(
        self, questions_list
    ) -> dict[str, int]:
        questions_answers = [question.answer for question in questions_list]
        answers_distribution = Counter(questions_answers)
        return answers_distribution

    def __get_subsection_average_distribution(
        self, subsection_questions
    ) -> dict[str, dict[str, int]]:
        subsection_distribution = self.__get_questions_average_distribution(
            questions_list=subsection_questions
        )
        return dict(subsection_distribution)

    def __get_section_and_subsection_answers_distribution(
        self, section
    ) -> dict[str, int]:
        subsections = section.answered_subsections.all()
        level3_questions = section.answered_questions_level_3.all()

        section_answers = level3_questions
        section_answers_with_comments = self.get_awnsers_with_comments(section_answers=section_answers)
        subsection_distribution = {}
        for subsction in subsections:
            subsection_title = subsction.form_subsection.subsection_title
            subsection_questions_level_1_and_level_2 = self.__get_subsection_questions(
                subsction
            )
            subsection_distribution[subsection_title] = (
                self.__get_subsection_average_distribution(
                    subsection_questions_level_1_and_level_2
                )
            )

            section_answers = section_answers.union(
                subsection_questions_level_1_and_level_2
            )

        section_distribution = self.__get_questions_average_distribution(
            section_answers
        )

        return section_distribution, subsection_distribution, section_answers_with_comments

    def __get_ona_form_total_metrics(
        self, distribution_by_section: dict
    ) -> dict[str, int]:
        total_counts = Counter()
        for section, counts in distribution_by_section.items():
            total_counts.update(counts)
        return dict(total_counts)

    def get_ona_form_average_distribution(
        self, ona_form
    ) -> dict[str, dict[str, int]]:
        sections = ona_form.answered_sections.all()

        distribution_by_section = {}
        distribution_by_subsections = {}
        for section in sections:
            section_name = section.form_section.section_title
            section_distribution, subsection_distribution, section_answers_with_comments = (
                self.__get_section_and_subsection_answers_distribution(section)
            )

            distribution_by_section[section_name] = dict(section_distribution)

            distribution_by_subsections[section_name] = dict(subsection_distribution)

        total_distribution = self.__get_ona_form_total_metrics(distribution_by_section)

        metrics = {
            "Subsections Distribution": distribution_by_subsections,
            "Sections Distribution": distribution_by_section,
            "ONA answer Distribution": total_distribution,
            "Answers with comments" : list(section_answers_with_comments)
        }

        return metrics
    
    def get_awnsers_with_comments(
            section_answers
    ):
        answers_with_comments = []
        for answerd_question in section_answers:
            if answerd_question.comment:
                print(answerd_question.comment)
                answers_with_comments.append(answerd_question.comment)
        return answers_with_comments
   
# Example usage:
if __name__ == "__main__":
    data = {'Subsections Distribution': {'2 SEÇÃO - ATENÇÃO AO PACIENTE': {'2.4 ATENDIMENTO CIRÚRGICO': {'parcial conforme': 14, 'conforme': 18, 'supera': 8, 'não conforme': 9}}}, 'Sections Distribution': {'2 SEÇÃO - ATENÇÃO AO PACIENTE': {'parcial conforme': 14, 'conforme': 18, 'supera': 8, 'não conforme': 9}}, 'ONA answer Distribution': {'parcial conforme': 14, 'conforme': 18, 'supera': 8, 'não conforme': 9}, 'Answers with comments': [{'id': 7, 'question': 266, 'answer': 'não conforme', 'comment': '\\sdf\\sdf\\sdfwecswdr4eysc dxz\\sdf\\sdf\\sdfwecswdr4eysc dxz\\sdf\\sdf\\sdfwecswdr4eysc dxz\\sdf\\sdf\\sdfwecswdr4eysc dxzVSDDSVCVCXVCX XCVFFB VCFD \\sdf\\sdf\\sdfwecswdr4eysc dxz'}, {'id': 1, 'question': 260, 'answer': 'não conforme', 'comment': 'dsfzsdzsdf'}, {'id': 25, 'question': 284, 'answer': 'não conforme', 'comment': 'vcbcvbbbcvbcvbc'}, {'id': 14, 'question': 273, 'answer': 'não conforme', 'comment': 'vbxcvbndeae'}, {'id': 6, 'question': 265, 'answer': 'não conforme', 'comment': '\\sdf\\sdf\\sdfwecswdr4eysc dxz'}, {'id': 41, 'question': 1380, 'answer': 'não conforme', 'comment': 'bzdfbzdfbxvc'}, {'id': 12, 'question': 271, 'answer': 'não conforme', 'comment': 'cvbdzfgreczxcvxcvcvzcfd'}, {'id': 31, 'question': 290, 'answer': 'não conforme', 'comment': 'cxvbcbcvbcvb'}, {'id': 21, 'question': 280, 'answer': 'não conforme', 'comment': 'cxzdfbrdhgercxcbcvbvc'}]}
    subsections = data["Subsections Distribution"]
    sections = data["Sections Distribution"]
    ona_awnser = data["ONA answer Distribution"]
    # graph = GraphsGenerator()
    # bar_plot_img = graph.plot_bar_plot(sections['SEÇÃO 01 - GESTÃO ORGANIZACIONAL'], 'SEÇÃO 01 - GESTÃO ORGANIZACIONAL')
    
    # rg = ReportGenerator()
    # prs = rg.make_report(data, "TEST06.pptx")
    # section01 = sections['SEÇÃO 01 - GESTÃO ORGANIZACIONAL']
    pdf = PDFReportGeneator(
        filename="test.pdf"
    )
    pdf.create_pdf_report_for_subsection(
        filename="TEST",
        evaluator_name="TESTAOZAO",
        answers=data
    )